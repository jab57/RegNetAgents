#!/usr/bin/env python3
"""
Generate RegNetAgents Conference Poster as PowerPoint
Reads content from REGNETAGENTS_CONFERENCE_POSTER.md and creates a professionally formatted PowerPoint poster
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re
import os

def parse_markdown_poster(filepath):
    """Parse the markdown conference poster file and extract structured content"""
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()

    sections = {}

    # Extract main title
    title_match = re.search(r'^# (.+)$', content, re.MULTILINE)
    sections['title'] = title_match.group(1) if title_match else "Gene Regulatory Agents"

    # Extract subtitle
    subtitle_match = re.search(r'^\*\*(.+?)\*\*$', content, re.MULTILINE)
    sections['subtitle'] = subtitle_match.group(1) if subtitle_match else ""

    # Extract abstract section
    abstract_match = re.search(r'## ABSTRACT\n\n(.+?)(?=\n\n##|\n---)', content, re.DOTALL)
    if abstract_match:
        sections['abstract'] = clean_text(abstract_match.group(1))

    # Extract innovation section
    innovation_match = re.search(r'## INNOVATION: Why This Is Novel\n\n### The Current State(.+?)\n\n### Our Innovation\n\n(.+?)\n\n### Quantitative Performance Comparison\n\n(.+?)\n\n### What Makes This Conference-Worthy\n\n(.+?)(?=\n\n##|\n---)', content, re.DOTALL)
    if innovation_match:
        sections['innovation_current'] = clean_text(innovation_match.group(1))
        sections['innovation_solution'] = extract_bullets(innovation_match.group(2))
        sections['innovation_table'] = parse_results_table(innovation_match.group(3))
        sections['innovation_worthy'] = extract_bullets(innovation_match.group(4))

    # Extract introduction
    intro_match = re.search(r'## INTRODUCTION\n\n### The Challenge\n(.+?)\n\n### Our Solution\n(.+?)(?=\n\n##|\n---)', content, re.DOTALL)
    if intro_match:
        sections['intro_challenge'] = extract_bullets(intro_match.group(1))
        sections['intro_solution'] = extract_bullets(intro_match.group(2))

    # Extract methods section
    methods_match = re.search(r'## METHODS\n\n### Multi-Agent Architecture\n\n```\n(.+?)\n```', content, re.DOTALL)
    if methods_match:
        sections['workflow_diagram'] = methods_match.group(1)

    # Extract specialized agents
    agents_match = re.search(r'### Specialized Agents\n\n(.+?)(?=\n### |## )', content, re.DOTALL)
    if agents_match:
        sections['agents'] = parse_agents_table(agents_match.group(1))

    # Extract data sources
    data_match = re.search(r'### Data Sources\n\n(.+?)(?=\n### |## )', content, re.DOTALL)
    if data_match:
        sections['data_sources'] = clean_text(data_match.group(1))

    # Extract natural language interface section
    nl_interface_match = re.search(r'## NATURAL LANGUAGE INTERFACE\n\n(.+?)(?=\n##|\n---)', content, re.DOTALL)
    if nl_interface_match:
        sections['nl_interface'] = clean_text(nl_interface_match.group(1))

    # Extract results - case study 1
    results_match = re.search(r'### Case Study 1: Multi-Gene Analysis - Rapid Biomarker Characterization for Colorectal Cancer Screening\n\n(.+?)(?=\n### Case Study 2:|\n### Performance)', content, re.DOTALL)
    if results_match:
        sections['biomarker_case'] = parse_biomarker_case(results_match.group(1))

    # Extract TP53 case study
    tp53_match = re.search(r'### Case Study 2: Single Gene Deep Dive - TP53 Analysis\n\n(.+?)(?=\n### Performance)', content, re.DOTALL)
    if tp53_match:
        sections['tp53_case'] = parse_tp53_case(tp53_match.group(1))

    # Extract performance metrics
    perf_match = re.search(r'### Performance Metrics\n\n(.+?)(?=\n##|\n---)', content, re.DOTALL)
    if perf_match:
        sections['performance'] = parse_performance_table(perf_match.group(1))

    # Extract data sources & validation section
    validation_match = re.search(r'## DATA SOURCES & VALIDATION\n\n(.+?)(?=\n##|\n---)', content, re.DOTALL)
    if validation_match:
        sections['validation'] = clean_text(validation_match.group(1))

    # Extract discussion
    discussion_match = re.search(r'## DISCUSSION\n\n### Key Advantages\n\n(.+?)\n\n### Limitations & Future Work\n\n(.+?)\n\n### Impact\n\n(.+?)(?=\n##|\n---)', content, re.DOTALL)
    if discussion_match:
        sections['advantages'] = extract_bullets(discussion_match.group(1))
        sections['limitations'] = extract_bullets(discussion_match.group(2))
        sections['impact'] = extract_bullets(discussion_match.group(3))

    # Extract conclusions
    conclusions_match = re.search(r'## CONCLUSIONS\n\n(.+?)(?=\n##|\n---)', content, re.DOTALL)
    if conclusions_match:
        sections['conclusions'] = clean_text(conclusions_match.group(1))

    # Extract take home message
    takehome_match = re.search(r'## TAKE HOME MESSAGE\n\n(.+?)(?=\n##|\n---)', content, re.DOTALL)
    if takehome_match:
        sections['takehome'] = clean_text(takehome_match.group(1))

    # Extract contact
    contact_match = re.search(r'## CONTACT\n\n(.+?)(?=\n##|\n---)', content, re.DOTALL)
    if contact_match:
        sections['contact'] = clean_text(contact_match.group(1))

    # Extract acknowledgments
    ack_match = re.search(r'## ACKNOWLEDGMENTS\n\n(.+?)(?=\n##|\n---)', content, re.DOTALL)
    if ack_match:
        sections['acknowledgments'] = clean_text(ack_match.group(1))

    return sections

def clean_text(text):
    """Clean markdown formatting from text"""
    # Remove markdown bold/italic
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    # Remove extra whitespace
    text = re.sub(r'\n\s*\n', '\n\n', text)
    return text.strip()

def extract_bullets(text):
    """Extract bullet points from text"""
    bullets = re.findall(r'^[-•]\s*(.+)$', text, re.MULTILINE)
    return [clean_text(b) for b in bullets]

def parse_agents_table(text):
    """Parse the agents table from markdown"""
    agents = []
    # Look for table rows
    rows = re.findall(r'\|\s*\*\*(.+?)\*\*\s*\|\s*(.+?)\s*\|\s*(.+?)\s*\|', text)
    for agent_name, function, output in rows:
        if agent_name not in ['Agent', '------']:  # Skip header/separator
            agents.append({
                'name': agent_name.strip(),
                'function': function.strip(),
                'output': output.strip()
            })
    return agents

def parse_biomarker_case(text):
    """Parse biomarker case study results"""
    case = {}

    # Extract query
    query_match = re.search(r'\*\*Query\*\*:\s*"(.+?)"', text)
    if query_match:
        case['query'] = query_match.group(1)

    # Extract analysis info
    analysis_match = re.search(r'\*\*Analysis\*\*:\s*(.+)', text)
    if analysis_match:
        case['analysis'] = analysis_match.group(1)

    # Extract execution time
    time_match = re.search(r'\*\*Execution Time\*\*:\s*(.+)', text)
    if time_match:
        case['time'] = time_match.group(1)

    # Extract table data
    table_match = re.search(r'\| Gene \|(.+?)(?=\n\n\*\*)', text, re.DOTALL)
    if table_match:
        case['results_table'] = parse_results_table(table_match.group(0))

    # Extract key findings
    findings_match = re.search(r'\*\*Key Findings\*\*:\n\n(.+?)(?=\n### |## )', text, re.DOTALL)
    if findings_match:
        case['findings'] = extract_bullets(findings_match.group(1))

    # Extract multi-agent insights
    insights_match = re.search(r'\*\*Multi-Agent Insights\*\*:\n(.+?)(?=\n\*\*Biomarker|\n### |## )', text, re.DOTALL)
    if insights_match:
        case['insights'] = extract_bullets(insights_match.group(1))

    # Extract biomarker type definitions
    definitions_match = re.search(r'\*\*Biomarker Type Definitions\*\*:\n(.+?)(?=\n\*\*|\n### |## )', text, re.DOTALL)
    if definitions_match:
        case['biomarker_definitions'] = extract_bullets(definitions_match.group(1))

    return case

def parse_tp53_case(text):
    """Parse TP53 case study"""
    case = {}

    # Extract execution time
    time_match = re.search(r'\*\*Execution Time\*\*:\s*(.+)', text)
    if time_match:
        case['time'] = time_match.group(1)

    # Extract network analysis
    network_match = re.search(r'\*\*Network Analysis\*\* \(RegNetAgentsModelingAgent\):\n(.+?)(?=\n\n\*\*)', text, re.DOTALL)
    if network_match:
        case['network'] = extract_bullets(network_match.group(1))

    # Extract perturbation analysis
    perturbation_match = re.search(r'\*\*Perturbation Analysis\*\*:\n(.+?)(?=\n\n\*\*)', text, re.DOTALL)
    if perturbation_match:
        case['perturbation'] = extract_bullets(perturbation_match.group(1))

    # Extract pathway enrichment
    pathway_match = re.search(r'\*\*Pathway Enrichment\*\* \(PathwayEnricherAgent\):\n(.+?)(?=\n\n\*\*)', text, re.DOTALL)
    if pathway_match:
        case['pathways'] = extract_bullets(pathway_match.group(1))

    # Extract domain analysis
    domain_match = re.search(r'\*\*Domain Analysis\*\*:\n(.+?)(?=\n\n\*\*)', text, re.DOTALL)
    if domain_match:
        case['domains'] = extract_bullets(domain_match.group(1))

    # Extract cross-cell
    cross_match = re.search(r'\*\*Cross-Cell Analysis\*\*:\s*(.+)', text)
    if cross_match:
        case['cross_cell'] = cross_match.group(1)

    return case

def parse_results_table(table_text):
    """Parse markdown table into structured data"""
    rows = []
    lines = table_text.strip().split('\n')

    for i, line in enumerate(lines):
        if '|' in line:
            cells = [c.strip() for c in line.split('|')[1:-1]]  # Remove empty first/last
            # Skip separator line (contains dashes)
            if cells and cells[0] and not cells[0].startswith('---'):
                # Clean markdown bold from gene names
                cells = [re.sub(r'\*\*(.+?)\*\*', r'\1', cell) for cell in cells]
                rows.append(cells)

    return rows

def parse_performance_table(text):
    """Parse performance metrics table"""
    metrics = {}
    rows = re.findall(r'\|\s*(.+?)\s*\|\s*(.+?)\s*\|', text)
    for metric, value in rows:
        if metric not in ['Metric', '------']:
            metrics[metric.strip()] = value.strip()
    return metrics

def create_poster(sections):
    """Create the PowerPoint poster from parsed sections"""

    # Get project root directory
    script_dir = os.path.dirname(os.path.abspath(__file__))
    project_root = os.path.dirname(script_dir)

    # Create presentation with poster dimensions (landscape: 48" x 36")
    prs = Presentation()
    prs.slide_width = Inches(48)
    prs.slide_height = Inches(36)

    # Add blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Define colors (RegNetAgents brand colors)
    DARK_BLUE = RGBColor(0, 51, 102)      # #003366
    LIGHT_BLUE = RGBColor(0, 102, 204)    # #0066CC
    GREEN = RGBColor(0, 170, 0)           # #00AA00
    ORANGE = RGBColor(255, 136, 0)        # #FF8800
    PURPLE = RGBColor(153, 51, 204)       # #9933CC
    GRAY = RGBColor(100, 100, 100)        # #646464
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)

    # ===== TITLE SECTION =====
    title_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.5),
        Inches(47), Inches(3.2)
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = DARK_BLUE
    title_box.line.color.rgb = DARK_BLUE

    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.margin_left = Inches(0)
    title_frame.margin_right = Inches(0)

    # Main title
    title = title_frame.add_paragraph()
    title.text = sections.get('title', 'Gene Regulatory Agents')
    title.font.size = Pt(72)
    title.font.bold = True
    title.font.color.rgb = WHITE
    title.alignment = PP_ALIGN.CENTER

    # Subtitle
    if sections.get('subtitle'):
        subtitle = title_frame.add_paragraph()
        subtitle.text = sections['subtitle']
        subtitle.font.size = Pt(40)
        subtitle.font.color.rgb = RGBColor(200, 230, 255)
        subtitle.alignment = PP_ALIGN.CENTER

    # Author and affiliation in separate textbox below title box
    # Title box ends at Y=3.7 (0.5 + 3.2), so start author at Y=3.8
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(47), Inches(0.8))
    author_frame = author_box.text_frame
    author_frame.word_wrap = False
    author_frame.margin_left = Inches(0)
    author_frame.margin_right = Inches(0)

    author_para = author_frame.add_paragraph()
    author_para.text = "Jose A. Bird, PhD  •  Bird AI Solutions"
    author_para.font.size = Pt(36)  # Slightly smaller to fit better
    author_para.font.bold = True
    author_para.font.color.rgb = WHITE
    author_para.alignment = PP_ALIGN.CENTER

    # ===== COLUMN 1: ABSTRACT & INNOVATION =====
    col1_x = Inches(0.5)
    col1_width = Inches(10)

    # Abstract - starts after author box (which ends at Y=4.6)
    # INCREASED height to accommodate long content (was 7.5", now 10.0")
    abstract_header = add_section_header(slide, col1_x, Inches(5.0), col1_width, "ABSTRACT", LIGHT_BLUE)

    abstract_box = slide.shapes.add_textbox(col1_x, Inches(6.0), col1_width, Inches(10.0))
    abstract_frame = abstract_box.text_frame
    abstract_frame.word_wrap = True
    abstract_frame.vertical_anchor = MSO_ANCHOR.TOP  # Align text to top

    if 'abstract' in sections:
        add_body_text(abstract_frame, sections['abstract'])

    # Innovation - Key Differentiator (moved down to accommodate larger abstract)
    # Abstract ends at 6.0 + 10.0 = 16.0, so start innovation at 16.5
    # INCREASED height to 6.0" to accommodate 8 bullets + performance stats
    innovation_header = add_section_header(slide, col1_x, Inches(16.5), col1_width, "WHY THIS IS NOVEL", ORANGE)

    innovation_box = slide.shapes.add_textbox(col1_x, Inches(17.5), col1_width, Inches(6.0))
    innovation_frame = innovation_box.text_frame
    innovation_frame.word_wrap = True

    if 'innovation_solution' in sections:
        add_subheader(innovation_frame, "Our Innovation")
        for bullet in sections['innovation_solution']:
            add_bullet(innovation_frame, bullet, bold=True)

    # Add comparison stats
    if 'innovation_table' in sections and len(sections['innovation_table']) > 1:
        add_subheader(innovation_frame, "\nPerformance Impact")
        # Show key comparisons
        add_bullet(innovation_frame, "Single gene (rule-based): Manual → 0.68 sec", bold=True)
        add_bullet(innovation_frame, "Single gene (LLM-powered): Manual → 15 sec", bold=True)
        add_bullet(innovation_frame, "5 genes (LLM-powered): Manual → 62 sec", bold=True)

    # Introduction - condensed
    # Innovation ends at 17.5 + 6.0 = 23.5, so start challenge at 24.0
    intro_header = add_section_header(slide, col1_x, Inches(24.0), col1_width, "THE CHALLENGE", LIGHT_BLUE)

    intro_box = slide.shapes.add_textbox(col1_x, Inches(25.0), col1_width, Inches(3.5))
    intro_frame = intro_box.text_frame
    intro_frame.word_wrap = True

    if 'intro_challenge' in sections:
        for bullet in sections['intro_challenge'][:5]:  # Limit to 5
            add_bullet(intro_frame, bullet)

    # Data Sources & Validation box - condensed to fit
    # Challenge ends at 25.0 + 3.5 = 28.5, so start data at 29.0
    data_header = add_section_header(slide, col1_x, Inches(29.0), col1_width, "DATA & VALIDATION", GREEN)

    data_box = slide.shapes.add_textbox(col1_x, Inches(30.0), col1_width, Inches(4.5))
    data_frame = data_box.text_frame
    data_frame.word_wrap = True

    add_subheader(data_frame, "Regulatory Networks")
    add_bullet(data_frame, "GREmLN foundation model (Zhang et al. 2025)")
    add_bullet(data_frame, "500K+ cells from CELLxGENE • ARACNe-AP algorithm")
    add_bullet(data_frame, "10 cell types: immune, blood, epithelial")

    add_subheader(data_frame, "\nPathway Enrichment")
    add_bullet(data_frame, "Reactome: Hypergeometric test, FDR < 0.05")

    add_subheader(data_frame, "\nValidation")
    add_bullet(data_frame, "CRC case: 5/5 genes literature-validated")
    add_bullet(data_frame, "All data from public sources")

    # ===== COLUMN 2: METHODS (WORKFLOW) =====
    col2_x = Inches(11)
    col2_width = Inches(11.5)

    methods_header = add_section_header(slide, col2_x, Inches(4.2), col2_width, "METHODS: Multi-Agent Architecture", LIGHT_BLUE)

    # Create simplified workflow diagram
    workflow_y = Inches(5.2)
    box_width = Inches(10.5)
    box_height = Inches(1.2)
    spacing = Inches(0.3)

    workflow_steps = [
        ("User Query (Claude)", DARK_BLUE, WHITE),
        ("Validate Inputs & Setup", LIGHT_BLUE, WHITE),
        ("RegNetAgentsModelingAgent:\nNetwork Analysis", GREEN, WHITE),
        ("RegNetAgentsModelingAgent:\nRegulators + Targets", GREEN, WHITE),
        ("RegNetAgentsModelingAgent:\nPerturbation Analysis", GREEN, WHITE),
        ("PathwayEnricherAgent:\nReactome Enrichment", PURPLE, WHITE),
    ]

    current_y = workflow_y
    for step_text, color, text_color in workflow_steps:
        add_workflow_box(slide, col2_x + Inches(0.5), current_y, box_width, Inches(1.2), step_text, color, text_color)
        current_y += box_height + spacing
        add_arrow(slide, col2_x + Inches(5.75), current_y - spacing + Inches(0.05), current_y - Inches(0.05), color)

    # Arrow to domain agents
    parallel_y = current_y + spacing
    add_arrow(slide, col2_x + Inches(5.75), current_y - spacing, parallel_y - Inches(0.1), GRAY)

    # Domain agents (PARALLEL - 4 simultaneously)
    domain_y = parallel_y
    domain_label = slide.shapes.add_textbox(col2_x + Inches(0.5), domain_y, box_width, Inches(0.6))
    domain_label_frame = domain_label.text_frame
    p = domain_label_frame.add_paragraph()
    p.text = "4 Domain Agents (PARALLEL via asyncio)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    domain_y += Inches(0.8)
    domain_box_width = Inches(2.5)

    domain_agents = ["Cancer\nAgent", "Drug\nAgent", "Clinical\nAgent", "Systems\nAgent"]
    domain_x = col2_x + Inches(0.5)
    for agent_text in domain_agents:
        add_workflow_box(slide, domain_x, domain_y, domain_box_width, Inches(1.2), agent_text, ORANGE, WHITE)
        domain_x += domain_box_width + Inches(0.1)

    # Final report
    final_y = domain_y + Inches(1.5)
    add_arrow(slide, col2_x + Inches(5.75), final_y, final_y + Inches(0.9), GRAY)
    add_workflow_box(slide, col2_x + Inches(0.5), final_y + Inches(1), box_width, box_height,
                     "Generate Comprehensive Report", PURPLE, WHITE)

    # Agent descriptions
    agent_desc_y = final_y + Inches(2.5)
    agent_header = add_section_header(slide, col2_x, agent_desc_y, col2_width, "Specialized Agents", DARK_BLUE)

    agent_box = slide.shapes.add_textbox(col2_x, agent_desc_y + Inches(1), col2_width, Inches(10))
    agent_frame = agent_box.text_frame
    agent_frame.word_wrap = True

    if 'agents' in sections:
        for agent in sections['agents']:
            p = agent_frame.add_paragraph()
            p.text = f"{agent['name']}: "
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE
            p.level = 0

            p = agent_frame.add_paragraph()
            p.text = f"{agent['function']} → {agent['output']}"
            p.font.size = Pt(18)
            p.font.color.rgb = BLACK
            p.level = 1
            p.space_after = Pt(8)

    # ===== KEY METRIC CALLOUT (top of column 3) =====
    col3_x = Inches(23)
    col3_width = Inches(12)

    # Prominent speed/validation callout
    add_callout_box(slide, col3_x, Inches(4.2), col3_width, Inches(1.2),
                   "⚡ 0.60-62 sec Analysis • 100% Literature-Validated ✓",
                   RGBColor(255, 250, 205), ORANGE)  # Light yellow bg, orange border

    # ===== COLUMN 3: NATURAL LANGUAGE INTERFACE & RESULTS =====
    # Natural Language Interface section (moved down to accommodate callout)
    nl_header = add_section_header(slide, col3_x, Inches(5.6), col3_width, "NATURAL LANGUAGE INTERFACE", PURPLE)

    nl_box = slide.shapes.add_textbox(col3_x, Inches(6.8), col3_width, Inches(4.5))
    nl_frame = nl_box.text_frame
    nl_frame.word_wrap = True

    if 'nl_interface' in sections:
        # Add main description
        add_body_text(nl_frame, "Conversational access via Claude Desktop - no programming required")

        # Add example prompts
        add_subheader(nl_frame, "\nExample Prompts:")
        add_bullet(nl_frame, '"Find biomarkers for CRC screening in Wnt/TP53 pathways"')
        add_bullet(nl_frame, '"Analyze TP53 - what does it regulate?"')
        add_bullet(nl_frame, '"Compare TP53, APC, BRCA1 across cell types"')

        add_body_text(nl_frame, "\nSystem auto-routes → second-scale response with comprehensive analysis")

    results_header = add_section_header(slide, col3_x, Inches(11.5), col3_width, "CASE 1: Multi-Gene Biomarker Analysis", LIGHT_BLUE)

    if 'biomarker_case' in sections:
        case = sections['biomarker_case']

        results_box = slide.shapes.add_textbox(col3_x, Inches(12.5), col3_width, Inches(2.5))
        results_frame = results_box.text_frame
        results_frame.word_wrap = True

        # Add user prompt
        p = results_frame.add_paragraph()
        p.text = "User Prompt:"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PURPLE

        p = results_frame.add_paragraph()
        p.text = '"Characterize these candidate genes for CRC biomarker potential: MYC, CTNNB1, CCND1, TP53, KRAS"'
        p.font.size = Pt(20)
        p.font.italic = True
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(8)

        # Add candidate selection explanation
        p = results_frame.add_paragraph()
        p.text = "Selected from known CRC pathways: Wnt (CTNNB1, MYC, CCND1) • TP53 pathway (TP53) • MAPK (KRAS)"
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.space_after = Pt(12)

        if 'query' in case:
            p = results_frame.add_paragraph()
            p.text = "Clinical Context: "
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE

            p = results_frame.add_paragraph()
            p.text = f'{case["query"]}'
            p.font.size = Pt(20)
            p.font.color.rgb = BLACK

        if 'analysis' in case:
            p = results_frame.add_paragraph()
            p.text = f"\n{case['analysis']}"
            p.font.size = Pt(22)
            p.font.color.rgb = BLACK

        if 'time' in case:
            p = results_frame.add_paragraph()
            p.text = f"Execution Time: {case['time']}"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = ORANGE

        # Results table
        if 'results_table' in case:
            table_top = Inches(15)
            create_results_table(slide, col3_x, table_top, col3_width, Inches(6.5), case['results_table'])

        # Biomarker type definitions (positioned after table ends at 21.5)
        if 'biomarker_definitions' in case:
            definitions_y = Inches(21.7)  # Table ends at 21.5, start definitions at 21.7
            definitions_box = slide.shapes.add_textbox(col3_x, definitions_y, col3_width, Inches(1.2))
            definitions_frame = definitions_box.text_frame
            definitions_frame.word_wrap = True

            p = definitions_frame.add_paragraph()
            p.text = "Biomarker Types:"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE

            for definition in case['biomarker_definitions']:
                p = definitions_frame.add_paragraph()
                p.text = definition
                p.font.size = Pt(16)
                p.font.color.rgb = BLACK
                p.space_after = Pt(2)

        # Multi-agent insights (positioned after definitions end at 22.9)
        if 'insights' in case:
            insights_y = Inches(23.0)  # Definitions end at 22.9, start insights at 23.0
            insights_header = add_section_header(slide, col3_x, insights_y, col3_width, "Multi-Agent Insights", ORANGE)

            insights_box = slide.shapes.add_textbox(col3_x, insights_y + Inches(1), col3_width, Inches(2.5))
            insights_frame = insights_box.text_frame
            insights_frame.word_wrap = True

            for insight in case['insights']:
                add_bullet(insights_frame, insight)

        # Key findings (condensed)
        # Insights box ends at 24.0 + 2.5 = 26.5, so start findings at 26.7
        if 'findings' in case:
            findings_y = Inches(26.7)
            findings_header = add_section_header(slide, col3_x, findings_y, col3_width, "Key Findings", GREEN)

            findings_box = slide.shapes.add_textbox(col3_x, findings_y + Inches(1), col3_width, Inches(2.5))
            findings_frame = findings_box.text_frame
            findings_frame.word_wrap = True

            for finding in case['findings'][:4]:  # Limit to 4 findings to save space
                add_bullet(findings_frame, finding, bold=True)

    # TP53 Case Study - CONDENSED VERSION showing perturbation analysis
    if 'tp53_case' in sections:
        tp53 = sections['tp53_case']
        tp53_y = Inches(24.5)  # Moved up to utilize available space
        tp53_header = add_section_header(slide, col3_x, tp53_y, col3_width, "CASE 2: TP53 Perturbation Analysis", PURPLE)

        # Plenty of room now: 36" - 25.5" = 10.5" available
        tp53_box = slide.shapes.add_textbox(col3_x, tp53_y + Inches(1), col3_width, Inches(6.5))
        tp53_frame = tp53_box.text_frame
        tp53_frame.word_wrap = True

        # Condensed version highlighting perturbation analysis
        p = tp53_frame.add_paragraph()
        p.text = "User Query: Comprehensive TP53 analysis with therapeutic target identification"
        p.font.size = Pt(20)
        p.font.italic = True
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(8)

        # Key network metrics
        p = tp53_frame.add_paragraph()
        p.text = "Network Analysis: Hub regulator • 163 targets • 7 regulators • 16 pathways (FDR<0.05)"
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)

        # Perturbation analysis - the key differentiator
        add_subheader(tp53_frame, "Automated Perturbation Analysis")

        p = tp53_frame.add_paragraph()
        p.text = "System simulated inhibition of all 7 upstream regulators and ranked therapeutic targets using standard network centrality metrics (PageRank, out-degree centrality):"
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.space_after = Pt(6)

        # Top targets with rankings
        add_bullet(tp53_frame, "Top Target (PageRank): WWTR1 (0.473) ✓ Validated Hippo pathway regulator")
        add_bullet(tp53_frame, "Top Target (Out-Degree): RBPMS (0.028, 403 targets) - Novel hypothesis")
        add_bullet(tp53_frame, "Validated: YAP1, CHD4 (Hippo pathway) confirmed in literature")
        add_bullet(tp53_frame, "Novel: PRRX2, THRA, IKZF2 prioritized for experimental validation")

        # Execution time
        p = tp53_frame.add_paragraph()
        p.text = "\nExecution: 0.60 sec (rule-based) | ~15 sec (LLM-powered with domain insights)"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ORANGE
        p.space_after = Pt(6)

        # Key insight
        p = tp53_frame.add_paragraph()
        p.text = "Demonstrates: Automated therapeutic target prioritization using PageRank (best predictor of drug target success per Mora & Donaldson 2021) identifies both validated regulators and novel hypotheses for experimental follow-up."
        p.font.size = Pt(17)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.space_after = Pt(4)

    # ===== COLUMN 4: DISCUSSION & CONCLUSIONS =====
    col4_x = Inches(35.5)
    col4_width = Inches(12)

    discussion_header = add_section_header(slide, col4_x, Inches(4.2), col4_width, "DISCUSSION", LIGHT_BLUE)

    discussion_box = slide.shapes.add_textbox(col4_x, Inches(5.2), col4_width, Inches(9))
    discussion_frame = discussion_box.text_frame
    discussion_frame.word_wrap = True
    discussion_frame.auto_size = None

    # Discussion content as paragraphs
    add_subheader(discussion_frame, "Key Advantages")

    p = discussion_frame.add_paragraph()
    p.text = "The system employs intelligent workflow routing that automatically selects optimal analysis paths based on gene characteristics, avoiding unnecessary computation while prioritizing relevant agents. Parallel agent execution provides 4-6x speedup versus sequential analysis through async I/O operations."
    p.font.size = Pt(20)
    p.font.color.rgb = BLACK
    p.space_after = Pt(10)

    p = discussion_frame.add_paragraph()
    p.text = "Domain-specific agents provide specialized expertise in cancer biology, drug development, clinical relevance, and systems biology with integrated scoring algorithms. Cell-type specificity across 10 pre-computed networks enables cross-cell comparison and tissue-specific pattern identification."
    p.font.size = Pt(20)
    p.font.color.rgb = BLACK
    p.space_after = Pt(10)

    add_subheader(discussion_frame, "\nLimitations & Future Work")

    p = discussion_frame.add_paragraph()
    p.text = "Current limitations include 10 cell types currently available (limited by available pre-computed networks) and networks updated periodically (not dynamic). Future work includes expanding cell type coverage as networks become available, integrating additional pathway databases (KEGG, GO, MSigDB), optimizing batch processing, and enhancing domain agent algorithms."
    p.font.size = Pt(20)
    p.font.color.rgb = BLACK
    p.space_after = Pt(10)

    add_subheader(discussion_frame, "\nAppropriate Use Cases")

    p = discussion_frame.add_paragraph()
    p.text = "System is designed for: (1) Hypothesis generation and target screening to prioritize experimental validation, (2) Network structure analysis to understand regulatory hierarchies and identify hub regulators, (3) Multi-gene comparative screening for biomarker prioritization, (4) Cross-cell-type pattern identification for tissue-specific regulation."
    p.font.size = Pt(20)
    p.font.color.rgb = BLACK
    p.space_after = Pt(10)

    p = discussion_frame.add_paragraph()
    p.text = "NOT designed for: Gene expression prediction (topology-based, not dynamical modeling), clinical diagnosis, or replacing experimental validation. Generates research hypotheses requiring wet-lab confirmation."
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_BLUE
    p.font.italic = True
    p.space_after = Pt(10)

    # Add Figure 4: Framework Value image (readable size)
    figure4_y = Inches(13.5)
    figure4_path = os.path.join(project_root, 'biorxiv', 'figure4_framework_value.png')

    if os.path.exists(figure4_path):
        try:
            # Add small caption above figure
            fig_caption_box = slide.shapes.add_textbox(col4_x, figure4_y, col4_width, Inches(0.25))
            fig_caption_frame = fig_caption_box.text_frame
            p = fig_caption_frame.add_paragraph()
            p.text = "Figure 4: Framework Performance (480-24,000× speedup)"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = ORANGE
            p.alignment = PP_ALIGN.CENTER

            # Add figure image (full column width for maximum readability)
            fig_pic = slide.shapes.add_picture(figure4_path, col4_x + Inches(0.2), figure4_y + Inches(0.3),
                                               width=col4_width - Inches(0.4))

            # Adjust conclusions header position to be below figure
            conclusions_y = figure4_y + Inches(0.3) + fig_pic.height + Inches(0.15)
        except Exception as e:
            print(f"   [WARNING] Could not add Figure 4: {e}")
            conclusions_y = Inches(16)
    else:
        conclusions_y = Inches(16)

    # Conclusions
    conclusions_header = add_section_header(slide, col4_x, conclusions_y, col4_width, "CONCLUSIONS", PURPLE)

    conclusions_box = slide.shapes.add_textbox(col4_x, conclusions_y + Inches(1), col4_width, Inches(3.2))
    conclusions_frame = conclusions_box.text_frame
    conclusions_frame.word_wrap = True
    conclusions_frame.auto_size = None

    # Conclusions content (compact)
    p = conclusions_frame.add_paragraph()
    p.text = "Gene Regulatory Agents demonstrates that multi-agent AI frameworks can effectively orchestrate complex gene regulatory network analyses across multiple biological domains. The framework successfully bridges computational biology, AI workflow orchestration, and conversational interfaces to make sophisticated gene analysis accessible to researchers."
    p.font.size = Pt(20)
    p.font.color.rgb = BLACK
    p.space_after = Pt(8)

    add_subheader(conclusions_frame, "Key Innovations")

    p = conclusions_frame.add_paragraph()
    p.text = "• Intelligent routing optimizes analysis paths based on gene characteristics"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.space_after = Pt(2)

    p = conclusions_frame.add_paragraph()
    p.text = "• Parallel execution enables second-scale analysis (0.6-62 sec depending on mode)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.space_after = Pt(2)

    p = conclusions_frame.add_paragraph()
    p.text = "• Domain-specific agents provide specialized cancer, drug, clinical, and systems perspectives"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.space_after = Pt(2)

    p = conclusions_frame.add_paragraph()
    p.text = "• Cell-type specificity supports translational research questions"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.space_after = Pt(8)

    # References section (positioned after conclusions)
    ref_y = conclusions_y + Inches(5.0)  # Conclusions box is 3.2" tall + 1.8" spacing for clearance
    ref_header = add_section_header(slide, col4_x, ref_y, col4_width, "REFERENCES", GRAY)

    ref_box = slide.shapes.add_textbox(col4_x, ref_y + Inches(1), col4_width, Inches(3))
    ref_frame = ref_box.text_frame
    ref_frame.word_wrap = True

    references = [
        "Zhang M, et al. (2025). GREmLN: A Cellular Regulatory Network-Aware Transcriptomics Foundation Model. bioRxiv 2025.07.03.663009.",
        "Margolin AA, et al. (2006). ARACNe algorithm. BMC Bioinformatics.",
        "Lachmann A, et al. (2016). ARACNe-AP: Adaptive Partitioning. Bioinformatics.",
        "Gillespie M, et al. (2022). Reactome pathway database. Nucleic Acids Res.",
        "CellxGene Data Portal. Chan Zuckerberg Initiative.",
        "LangGraph Framework. LangChain AI.",
        "Model Context Protocol (MCP). Anthropic."
    ]

    for i, ref in enumerate(references, 1):
        p = ref_frame.add_paragraph()
        p.text = f"{i}. {ref}"
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK
        p.space_after = Pt(4)

    # ===== CONTACT & QR CODE (Bottom - Column 4) =====
    contact_y = ref_y + Inches(4)  # References section + spacing
    contact_x = Inches(35.5)
    contact_width = Inches(12)

    if 'contact' in sections:
        contact_header = add_section_header(slide, contact_x, contact_y, contact_width, "CONTACT", PURPLE)

        contact_box = slide.shapes.add_textbox(contact_x, contact_y + Inches(1), Inches(7.5), Inches(3.5))
        contact_frame = contact_box.text_frame
        contact_frame.word_wrap = True

        p = contact_frame.add_paragraph()
        p.text = "Jose A. Bird, PhD"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        p = contact_frame.add_paragraph()
        p.text = "Bird AI Solutions"
        p.font.size = Pt(20)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)

        p = contact_frame.add_paragraph()
        p.text = "Email: jbird@birdaisolutions.com"
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK

        p = contact_frame.add_paragraph()
        p.text = "LinkedIn: /in/jose-bird-data-science-advanced-analytics"
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)

        p = contact_frame.add_paragraph()
        p.text = "Interested in Gene Regulatory Agents?"
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = GRAY

        # QR Code placeholder
        qr_x = contact_x + Inches(8)
        qr_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, qr_x, contact_y + Inches(1),
                                        Inches(3.5), Inches(3.5))
        qr_box.fill.solid()
        qr_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
        qr_box.line.color.rgb = BLACK
        qr_box.line.width = Pt(2)

        qr_frame = qr_box.text_frame
        qr_frame.vertical_anchor = 1  # Middle
        p = qr_frame.add_paragraph()
        p.text = "QR CODE\nPoster\nWebpage"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(34.5), Inches(47), Inches(1))
    footer_frame = footer_box.text_frame
    footer_frame.word_wrap = True

    p = footer_frame.add_paragraph()
    p.text = "Powered by LangGraph | MCP Protocol | Claude Desktop"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    return prs

# Helper functions
def add_section_header(slide, x, y, width, text, color):
    """Add a colored section header"""
    header_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, Inches(0.8))
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = color
    header_box.line.color.rgb = color

    header_frame = header_box.text_frame
    p = header_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    return header_box

def add_workflow_box(slide, x, y, width, height, text, bg_color, text_color):
    """Add a workflow process box"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = RGBColor(100, 100, 100)
    box.line.width = Pt(2)

    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = 1  # Middle

    p = frame.add_paragraph()
    p.text = text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    return box

def add_arrow(slide, x, y_start, y_end, color):
    """Add a vertical arrow between boxes"""
    from pptx.enum.shapes import MSO_CONNECTOR
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y_start, x, y_end)
    connector.line.color.rgb = color
    connector.line.width = Pt(3)

def add_body_text(text_frame, text):
    """Add body text paragraph"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_after = Pt(12)
    return p

def add_subheader(text_frame, text):
    """Add subheader text"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_before = Pt(8)
    p.space_after = Pt(6)
    return p

def add_bullet(text_frame, text, bold=False):
    """Add bullet point"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(20)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.level = 1
    p.space_after = Pt(4)
    return p

def add_callout_box(slide, x, y, width, height, text, bg_color, border_color):
    """Add a prominent callout box for key metrics"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = border_color
    box.line.width = Pt(4)

    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = 1  # Middle
    frame.margin_left = Inches(0.2)
    frame.margin_right = Inches(0.2)

    p = frame.add_paragraph()
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    return box

def create_results_table(slide, x, y, width, height, data_rows):
    """Create results table with biomarker data and color-coded therapeutic scores"""
    # data_rows[0] contains the headers from markdown
    rows = len(data_rows)
    cols = len(data_rows[0]) if data_rows else 5

    table = slide.shapes.add_table(rows, cols, x, y, width, height).table

    # Set column widths (must be integer)
    col_width = int(width / cols)
    for i in range(cols):
        table.columns[i].width = col_width

    # Find the "Therapeutic Score" column index
    therapeutic_col_idx = -1
    if data_rows and len(data_rows) > 0:
        for idx, header in enumerate(data_rows[0]):
            if "Therapeutic" in header or "Score" in header:
                therapeutic_col_idx = idx
                break

    # Process all rows (first row is header from markdown)
    for row_idx, row_data in enumerate(data_rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text

            # First row is header
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(22)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Color code therapeutic score column
                if col_idx == therapeutic_col_idx and therapeutic_col_idx != -1:
                    try:
                        score = float(cell_text)
                        cell.fill.solid()
                        if score >= 0.75:
                            cell.fill.fore_color.rgb = RGBColor(144, 238, 144)  # Light green
                        elif score >= 0.40:
                            cell.fill.fore_color.rgb = RGBColor(255, 255, 153)  # Light yellow
                        else:
                            cell.fill.fore_color.rgb = RGBColor(255, 182, 193)  # Light red
                    except ValueError:
                        # Not a number, use default coloring
                        if row_idx % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                else:
                    # Data rows - alternating colors for other columns
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(240, 240, 240)

                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = RGBColor(0, 0, 0)

                # Bold first column (gene names) and therapeutic score column
                if col_idx == 0 or col_idx == therapeutic_col_idx:
                    paragraph.font.bold = True

                # Center align therapeutic scores
                if col_idx == therapeutic_col_idx:
                    paragraph.alignment = PP_ALIGN.CENTER

    return table

def main():
    """Main function to generate the poster"""
    print("=" * 60)
    print("RegNetAgents Conference Poster Generator")
    print("=" * 60)

    # Get project root directory
    script_dir = os.path.dirname(os.path.abspath(__file__))
    project_root = os.path.dirname(script_dir)

    # Parse markdown file
    print("\n1. Parsing REGNETAGENTS_CONFERENCE_POSTER.md...")
    try:
        # Try docs directory first, then current directory
        poster_path = os.path.join(project_root, 'docs', 'REGNETAGENTS_CONFERENCE_POSTER.md')
        if not os.path.exists(poster_path):
            poster_path = 'REGNETAGENTS_CONFERENCE_POSTER.md'
        sections = parse_markdown_poster(poster_path)
        print(f"   [OK] Extracted {len(sections)} sections from {poster_path}")
    except FileNotFoundError:
        print("   [ERROR] REGNETAGENTS_CONFERENCE_POSTER.md not found")
        print("   Make sure the file exists in docs/ directory or current directory")
        return
    except Exception as e:
        print(f"   [ERROR] parsing markdown: {e}")
        return

    # Create PowerPoint
    print("\n2. Generating PowerPoint presentation...")
    try:
        prs = create_poster(sections)
        print("   [OK] Poster layout created")
    except Exception as e:
        print(f"   [ERROR] creating poster: {e}")
        import traceback
        traceback.print_exc()
        return

    # Save file to biorxiv directory
    biorxiv_dir = os.path.join(project_root, "biorxiv")
    os.makedirs(biorxiv_dir, exist_ok=True)
    output_file = os.path.join(biorxiv_dir, "REGNETAGENTS_CONFERENCE_POSTER.pptx")
    print(f"\n3. Saving to {output_file}...")
    try:
        prs.save(output_file)
        print(f"   [OK] File saved successfully")
    except Exception as e:
        print(f"   [ERROR] saving file: {e}")
        return

    print("\n" + "=" * 60)
    print("SUCCESS! Conference poster generated")
    print("=" * 60)
    print(f"\nFile: {output_file}")
    print(f"Dimensions: 48\" × 36\" (landscape)")
    print(f"\nNext steps:")
    print(f"  1. Open {output_file} in PowerPoint")
    print(f"  2. Add institution logos or QR codes")
    print(f"  3. Review and adjust spacing if needed")
    print(f"  4. Export as PDF for printing (File -> Save As -> PDF)")
    print(f"\nReady for conference poster printing services!")
    print("=" * 60)

if __name__ == "__main__":
    try:
        main()
    except ImportError as e:
        print("\n" + "=" * 60)
        print("ERROR: python-pptx library not found")
        print("=" * 60)
        print("\nInstall with:")
        print("  pip install python-pptx")
        print(f"\nDetails: {e}")
    except Exception as e:
        print(f"\nERROR: {e}")
        import traceback
        traceback.print_exc()
